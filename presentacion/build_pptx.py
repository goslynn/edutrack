#!/usr/bin/env python3
"""Genera presentacion/edutrack.pptx (versión institucional, simple)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

BASE = "/home/vgonz/code/edutrack/presentacion"
DIAG = os.path.join(BASE, "assets", "diagramas")

INK   = RGBColor(0x0F, 0x17, 0x2A)
BLUE  = RGBColor(0x25, 0x63, 0xEB)
BLUED = RGBColor(0x1E, 0x40, 0xAF)
MUT   = RGBColor(0x64, 0x74, 0x8B)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
BG    = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x0F, 0x17, 0x2A)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def box(s, l, t, w, h):
    return s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def settext(tf, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
            space=6, font="Calibri"):
    """runs: str o lista de párrafos. Cada párrafo str o lista de (txt,bold,color)."""
    if isinstance(runs, str):
        runs = [runs]
    tf.word_wrap = True
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space)
        segs = para if isinstance(para, list) else [(para, bold, color)]
        for seg in segs:
            txt, b, c = seg
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = b
            r.font.color.rgb = c; r.font.name = font
    return tf


def eyebrow(s, txt, color=BLUE):
    tb = box(s, 0.7, 0.45, 11, 0.5)
    settext(tb.text_frame, [[(txt.upper(), True, color)]], size=13)


def title(s, txt, color=INK):
    tb = box(s, 0.7, 0.85, 12, 1.0)
    settext(tb.text_frame, [[(txt, True, color)]], size=34)


def bullets(s, l, t, w, h, items, size=15, color=INK):
    tb = box(s, l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(7)
        segs = it if isinstance(it, list) else [(it, False, color)]
        r0 = p.add_run(); r0.text = "•  "
        r0.font.size = Pt(size); r0.font.color.rgb = BLUE; r0.font.bold = True
        for txt, b, c in segs:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = b; r.font.color.rgb = c
    return tb


def card(s, l, t, w, h, fill=BG):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0); sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def foot(s):
    return  # sin pie de página ("EduTrack · DSY1106" eliminado por pedido)


# ---------- 1 · Portada ----------
s = slide(); bg(s, DARK)
tb = box(s, 0.9, 2.1, 11.5, 0.5)
settext(tb.text_frame, [[("COLEGIO BERNARDO O'HIGGINS · COQUIMBO", True, RGBColor(0x60,0xA5,0xFA))]], size=15)
tb = box(s, 0.85, 2.55, 11.5, 1.4)
settext(tb.text_frame, [[("EduTrack", True, WHITE)]], size=66)
tb = box(s, 0.9, 3.95, 11.5, 0.7)
settext(tb.text_frame, [[("Libro de Clases Digital — arquitectura de microservicios", False, RGBColor(0xCB,0xD5,0xE1))]], size=22)
tb = box(s, 0.9, 5.2, 11.5, 1.2)
settext(tb.text_frame, [
    [("DSY1106 · Desarrollo Fullstack III  |  Evaluación Parcial N°1", False, RGBColor(0x94,0xA3,0xB8))],
    [("Integrantes: {{ COMPLETAR NOMBRES }}   ·   {{ FECHA }}", False, RGBColor(0xCB,0xD5,0xE1))],
], size=15, space=8)

# ---------- 2 · Contexto y objetivos ----------
s = slide(); bg(s, WHITE)
eyebrow(s, "01 · Contexto y objetivos"); title(s, "El problema del colegio")
bullets(s, 0.7, 1.9, 6.1, 3, [
    [("Libro de clases ", False, INK), ("físico", True, BLUED), (" + soluciones digitales ", False, INK), ("fragmentadas", True, BLUED), (".", False, INK)],
    "Información dispersa, sin historial consolidado ni acceso oportuno.",
    [("Datos de ", False, INK), ("menores de edad", True, BLUED), (" → seguridad y trazabilidad (Ley 19.628).", False, INK)],
    [("EduTrack: entorno seguro para ", False, INK), ("docentes, administradores y apoderados.", True, BLUED)],
], size=16)
card(s, 7.1, 1.9, 5.4, 4.4)
tb = box(s, 7.4, 2.05, 4.9, 0.6)
settext(tb.text_frame, [[("Objetivos funcionales", True, BLUED)]], size=17)
bullets(s, 7.4, 2.6, 4.9, 3.6, [
    "Auth con roles dinámicos y permisos configurables",
    "Cursos con acceso granular por docente",
    "Alumnos y apoderados (avisos por correo)",
    "Notas y asistencia por clase",
    "Anotaciones con notificación a apoderados",
    "Reportes JSON / CSV / PDF",
], size=14)
foot(s)

# ---------- 3 · Plan vs realidad (tabla) ----------
s = slide(); bg(s, WHITE)
eyebrow(s, "02 · Lo que planteamos vs. dónde estamos"); title(s, "Cambios en el camino")
rows = [
    ("Tema", "Planteamos (entrega 1)", "Dónde estamos hoy"),
    ("Orquestación", "Kubernetes / Amazon EKS + HPA", "Fly.io — cada MS es una app"),
    ("Service discovery", "Consul + DNS de K8s", "Contrato de nombres vía *.fly.internal"),
    ("API Gateway", "Gateway genérico", "OpenResty + Lua (jwt.lua)"),
    ("Base de datos", "Amazon RDS", "PostgreSQL — shared DB + schema por MS"),
    ("Permisos", "(role, resource_uuid)", "resource_key string estable <svc>.<rec>"),
    ("Capa extra", "—", "BFF (TypeScript + Hono), no previsto"),
    ("S3 · Notif · Report · Contenido", "en el diseño", "aún no implementados"),
]
tw, th = Inches(11.9), Inches(4.5)
gt = s.shapes.add_table(len(rows), 3, Inches(0.7), Inches(1.95), tw, th).table
gt.columns[0].width = Inches(3.0); gt.columns[1].width = Inches(4.2); gt.columns[2].width = Inches(4.7)
for ci in range(3):
    c = gt.cell(0, ci); c.fill.solid(); c.fill.fore_color.rgb = BLUE
    c.text_frame.paragraphs[0].runs and None
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        c = gt.cell(ri, ci)
        c.margin_left = Inches(0.1); c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = c.text_frame; tf.word_wrap = True; tf.clear()
        p = tf.paragraphs[0]; r = p.add_run(); r.text = val
        r.font.size = Pt(12.5); r.font.name = "Calibri"
        if ri == 0:
            r.font.bold = True; r.font.color.rgb = WHITE
        else:
            r.font.color.rgb = INK
            if not (ri % 2):
                c.fill.solid(); c.fill.fore_color.rgb = BG
            else:
                c.fill.solid(); c.fill.fore_color.rgb = WHITE
foot(s)

# ---------- 4 · Estado: 6/9 ----------
s = slide(); bg(s, WHITE)
eyebrow(s, "02 · Estado actual"); title(s, "Microservicios: 6 de 9 construidos")
card(s, 0.7, 2.0, 5.8, 3.3)
tb = box(s, 1.0, 2.15, 5.3, 0.6); settext(tb.text_frame, [[("✅ Implementados y probados", True, GREEN)]], size=17)
bullets(s, 1.0, 2.75, 5.3, 2.2, [
    "auth · course · student",
    "attendance · assessment · annotation",
    [("+ Gateway (OpenResty), BFF (Hono), front (React 19)", False, MUT)],
], size=15)
card(s, 6.9, 2.0, 5.6, 3.3)
tb = box(s, 7.2, 2.15, 5.1, 0.6); settext(tb.text_frame, [[("⏳ Pendientes", True, MUT)]], size=17)
bullets(s, 7.2, 2.75, 5.1, 2.0, [
    "content · notification · report",
    [("Servicios satélite: contenido/S3, notificaciones async, reportes", False, MUT)],
], size=15)
tb = box(s, 0.7, 5.6, 11.9, 1.0)
settext(tb.text_frame, [[("Lectura honesta: ", True, BLUED), ("el núcleo de seguridad y dominio académico está construido y probado; faltan los servicios satélite.", False, INK)]], size=15)
foot(s)

# ---------- 5 · Stack ----------
s = slide(); bg(s, WHITE)
eyebrow(s, "03 · Stack final e infraestructura"); title(s, "Cómo está construido")
bullets(s, 0.7, 1.95, 11.8, 4.5, [
    [("Backend: ", True, BLUED), ("Quarkus 3.x + Java 21, Panache (Active Record), Flyway, Bean Validation", False, INK)],
    [("Gateway: ", True, BLUED), ("OpenResty + Lua — valida JWT RS256, propaga X-User-Id / X-User-Roles", False, INK)],
    [("BFF: ", True, BLUED), ("TypeScript + Hono", False, INK)],
    [("Frontend: ", True, BLUED), ("React 19 + Vite + Tailwind + shadcn", False, INK)],
    [("Datos: ", True, BLUED), ("PostgreSQL — shared DB, schema y credenciales por MS", False, INK)],
    [("Resiliencia: ", True, BLUED), ("SmallRye Fault Tolerance (Circuit Breaker)", False, INK)],
    [("Hosting: ", True, BLUED), ("Fly.io · *.fly.internal · local: Docker Compose", False, INK)],
], size=17)
foot(s)

# ---------- 6 · Diagrama arquitectura ----------
s = slide(); bg(s, WHITE)
eyebrow(s, "03 · Diagrama de arquitectura")
s.shapes.add_picture(os.path.join(DIAG, "arquitectura.png"), Inches(0.6), Inches(1.05),
                     width=Inches(12.1))

# ---------- 7 · Caso destacado intro ----------
s = slide(); bg(s, WHITE)
eyebrow(s, "04 · Caso destacado · patrón de diseño")
title(s, "Toma de asistencia: principio Abierto/Cerrado")
tb = box(s, 0.7, 1.85, 11.8, 0.7)
settext(tb.text_frame, [[("Contrato de asistencia agnóstico al método de captura: web hoy, biometría/RFID mañana, sin reescribir el servicio.", False, INK)]], size=16)
card(s, 0.7, 2.7, 11.8, 2.6)
tb = box(s, 1.0, 2.85, 11.2, 0.5); settext(tb.text_frame, [[("El patrón en una frase (Open/Closed Principle)", True, BLUED)]], size=17)
bullets(s, 1.0, 3.45, 11.2, 1.7, [
    [("Abierto a extensión: ", True, BLUED), ("un método de captura nuevo = un Adapter externo nuevo.", False, INK)],
    [("Cerrado a modificación: ", True, BLUED), ("el microservicio de asistencia nunca cambia su código para soportarlo.", False, INK)],
], size=16)
foot(s)

# ---------- 8 · Diagrama toma de asistencia ----------
s = slide(); bg(s, WHITE)
eyebrow(s, "04 · Flujo de la toma de asistencia (Abierto/Cerrado)")
s.shapes.add_picture(os.path.join(DIAG, "toma-asistencia.png"), Inches(0.6), Inches(1.05),
                     width=Inches(12.1))

# ---------- 9 · Patrones ----------
s = slide(); bg(s, WHITE)
eyebrow(s, "04 · Patrones de diseño en juego"); title(s, "Qué patrones lo sostienen")
prow = [
    ("Patrón / principio", "Cómo aparece aquí"),
    ("Open/Closed Principle (SOLID)", "nuevo método de captura sin modificar el microservicio"),
    ("Adapter", "cada origen (biometría, RFID...) traduce su dato crudo al contrato normalizado"),
    ("DTO + Mapper", "CreateRecordRequest -> dominio -> RecordResponse (capas desacopladas)"),
    ("Guarda de estado (State)", "una sesión CLOSED rechaza nuevos registros (validateOpen())"),
    ("Repository", "acceso a datos detrás de AttendanceRecordRepository (Panache)"),
]
gt = s.shapes.add_table(len(prow), 2, Inches(0.7), Inches(1.95), Inches(11.9), Inches(4.0)).table
gt.columns[0].width = Inches(4.9); gt.columns[1].width = Inches(7.0)
for ri, row in enumerate(prow):
    for ci, val in enumerate(row):
        c = gt.cell(ri, ci); c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.12); c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
        tf = c.text_frame; tf.word_wrap = True; tf.clear()
        p = tf.paragraphs[0]; r = p.add_run(); r.text = val
        r.font.size = Pt(13); r.font.name = "Calibri"
        if ri == 0:
            r.font.bold = True; r.font.color.rgb = WHITE
            c.fill.solid(); c.fill.fore_color.rgb = BLUE
        else:
            r.font.color.rgb = INK
            c.fill.solid(); c.fill.fore_color.rgb = (BG if ri % 2 == 0 else WHITE)
            if ci == 0:
                r.font.bold = True; r.font.color.rgb = BLUED
foot(s)

# ---------- 10 · Testing resultados ----------
s = slide(); bg(s, WHITE)
eyebrow(s, "05 · Testing unitario — resultados reales")
title(s, "206 tests backend · ~96% cobertura front")
# Backend table
bt = [("Servicio","Tests"),("auth","62"),("student","40"),("attendance","38"),
      ("assessment","33"),("course","17"),("annotation","16"),("TOTAL","206 · 0 fallos")]
g1 = s.shapes.add_table(len(bt), 2, Inches(0.7), Inches(2.0), Inches(5.5), Inches(4.2)).table
g1.columns[0].width = Inches(3.7); g1.columns[1].width = Inches(1.8)
for ri,row in enumerate(bt):
    for ci,val in enumerate(row):
        c=g1.cell(ri,ci); c.vertical_anchor=MSO_ANCHOR.MIDDLE; c.margin_left=Inches(0.12)
        tf=c.text_frame; tf.clear(); p=tf.paragraphs[0]; r=p.add_run(); r.text=val
        r.font.size=Pt(14); r.font.name="Calibri"
        last = ri==len(bt)-1
        if ri==0:
            r.font.bold=True; r.font.color.rgb=WHITE; c.fill.solid(); c.fill.fore_color.rgb=BLUE
        else:
            r.font.color.rgb=INK; r.font.bold=last
            c.fill.solid(); c.fill.fore_color.rgb=(RGBColor(0xDC,0xFC,0xE7) if last else (BG if ri%2==0 else WHITE))
# Frontend coverage card
card(s, 6.9, 2.0, 5.6, 4.2)
tb=box(s,7.2,2.15,5.0,0.5); settext(tb.text_frame,[[("Frontend — cobertura (Vitest)", True, GREEN)]], size=16)
bullets(s,7.2,2.8,5.0,3.0,[
    [("Statements   ", True, INK),("96.0%  (875/911)", False, INK)],
    [("Branches     ", True, INK),("88.8%  (411/463)", False, INK)],
    [("Functions    ", True, INK),("97.4%  (263/270)", False, INK)],
    [("Lines        ", True, INK),("97.6%  (772/791)", False, INK)],
    [("26 archivos de test · front/coverage/index.html", False, MUT)],
], size=15)
foot(s)

# ---------- 11 · Demo ----------
s = slide(); bg(s, WHITE)
eyebrow(s, "06 · Demo en vivo"); title(s, "Guion de la demo")
bullets(s, 0.7, 2.1, 11.8, 4.2, [
    [("1.  ", True, BLUE),("docker compose up → abrir el frontend", False, INK)],
    [("2.  ", True, BLUE),("Login y manejo de sesión (JWT)", False, INK)],
    [("3.  ", True, BLUE),("Abrir una sesión de clase y tomar asistencia (caso destacado en vivo)", False, INK)],
    [("4.  ", True, BLUE),("Registrar nota / anotación con validación en tiempo real", False, INK)],
    [("5.  ", True, BLUE),("Dashboard con estadísticas", False, INK)],
], size=19)
tb=box(s,0.7,6.0,11.8,0.5)
settext(tb.text_frame,[[("Páginas: dashboard · estudiantes · calificaciones · asistencia · anotaciones.", False, MUT)]], size=13)
foot(s)

# ---------- 12 · Cierre ----------
s = slide(); bg(s, DARK)
tb=box(s,0.9,1.4,11.5,0.5); settext(tb.text_frame,[[("CIERRE", True, RGBColor(0x60,0xA5,0xFA))]], size=15)
tb=box(s,0.85,1.85,11.5,1.0); settext(tb.text_frame,[[("En una frase", True, WHITE)]], size=34)
bullets(s,0.9,3.0,11.5,2.6,[
    [("Núcleo de seguridad y dominio académico construido y probado (206 tests).", False, RGBColor(0xE2,0xE8,0xF0))],
    [("Arquitectura de microservicios desplegable en Fly.io, simple y sin Consul.", False, RGBColor(0xE2,0xE8,0xF0))],
    [("Pendientes claros: content, notification, report.", False, RGBColor(0xE2,0xE8,0xF0))],
], size=18)
tb=box(s,0.9,5.7,11.5,0.9); settext(tb.text_frame,[[("Gracias — ¿preguntas?", True, WHITE)]], size=30)

out = os.path.join(BASE, "edutrack.pptx")
prs.save(out)
print("OK ->", out, f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
